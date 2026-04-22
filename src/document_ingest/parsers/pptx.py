# src/document_ingest/parsers/pptx.py
"""PPTX 파서 — python-pptx 기반 (M_06 스펙 §5.3)."""

from __future__ import annotations

import logging
from pathlib import Path

from document_ingest.errors import ParseError
from document_ingest.segments import _Segment

logger = logging.getLogger(__name__)


def _parse_pptx(path: str) -> list[_Segment]:
    """PPTX 파일에서 슬라이드 단위 텍스트 세그먼트를 추출한다.

    - page: 슬라이드 번호 (1-based).
    - section: 해당 슬라이드 제목 (없으면 None).
    - 노트(speaker notes): V1 skip.

    Args:
        path: 절대 또는 상대 경로 문자열.

    Returns:
        list[_Segment]

    Raises:
        ParseError: python-pptx 로드 실패.
    """
    try:
        import pptx
    except ImportError as exc:
        raise ParseError(f"python-pptx 패키지가 설치되지 않았습니다: {exc}") from exc

    try:
        prs = pptx.Presentation(path)
    except Exception as exc:
        raise ParseError(f"PPTX 열기 실패: {Path(path).name}: {exc}") from exc

    segments: list[_Segment] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title: str | None = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip()
            slide_title = title_text or None

        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)

        combined = "\n".join(texts)
        if combined.strip():
            segments.append(
                _Segment(
                    text=combined,
                    page=slide_idx,
                    section=slide_title,
                    bbox=None,
                )
            )

    return segments
